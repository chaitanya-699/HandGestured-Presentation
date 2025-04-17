import cv2
import numpy as np
import fitz
from PIL import Image
import base64
from io import BytesIO
import mediapipe as mp
from collections import deque
import time
from typing import Any
import logging
from pptx import Presentation
import os
import tempfile
import subprocess
from pdf2image import convert_from_path

class HandGestureProcessor:
    def __init__(self):
        self.mp_hands = mp.solutions.hands
        self.hands = self.mp_hands.Hands(
            static_image_mode=False,
            max_num_hands=1,
            min_detection_confidence=0.8,
            min_tracking_confidence=0.8
        )

        # Screen dimensions
        self.SLIDE_WIDTH = 1280
        self.SLIDE_HEIGHT = 720

        # Drawing settings
        self.colors = {
            'red': (0, 0, 255),
            'green': (0, 255, 0),
            'blue': (255, 0, 0),
            'yellow': (0, 255, 255),
            'purple': (255, 0, 255),
            'orange': (0, 165, 255)
        }
        
        # Drawing styles
        self.current_brush_style = "normal"
        self.current_color = self.colors['red']
        self.brush_thickness = 3
        self.pointer_color = (0, 255, 255)
        self.current_tool = "brush"
        
        # UI Control dimensions
        self.CONTROL_HEIGHT = 40
        self.BUTTON_WIDTH = 60
        self.BUTTON_SPACING = 5
        
        # Setup UI elements
        self.ui_elements = {}
        self.setup_ui()
        
        # Drawing state
        self.prev_x = None
        self.prev_y = None
        self.drawing_mode = False
        self.shape_start = None
        
        # Presentation state
        self.current_slide = 0
        self.slides = deque(maxlen=100)
        self.drawings = deque(maxlen=100)
        self.whiteboard = np.ones((self.SLIDE_HEIGHT, self.SLIDE_WIDTH, 3), dtype=np.uint8) * 255
        self.whiteboard_drawing = np.zeros((self.SLIDE_HEIGHT, self.SLIDE_WIDTH, 3), dtype=np.uint8)
        
        # Performance tracking
        self.fps_history = deque(maxlen=30)
        self.last_frame_time = time.time()
        
        # Gesture handling
        self.position_history = deque(maxlen=5)
        self.last_gesture = None
        self.gesture_cooldown = 0
        self.hover_cooldown = 0
        self.gesture_active = False
        self.navigation_timer = 0
        
        # Mode tracking
        self.is_whiteboard = False
        self.is_eraser = False
        
        self.hover_states = {
            'eraser': {
                'is_hovering': False,
                'cooldown': 0
            },
            'whiteboard': {
                'is_hovering': False,
                'cooldown': 0
            }
        }
        self.HOVER_COOLDOWN = 30 
        
        # Add caching for performance
        self._frame_cache = {}
        self._gesture_cache = {}
        self.CACHE_SIZE = 30
        
        # Pre-calculate commonly used values
        self._control_panel_rect = (100, 0, self.SLIDE_WIDTH - 100, self.CONTROL_HEIGHT)
        self._brush_styles = ["normal", "spray", "calligraphy", "neon"]
        
        # Use numpy arrays for better performance
        self._zero_drawing = np.zeros((self.SLIDE_HEIGHT, self.SLIDE_WIDTH, 3), dtype=np.uint8)
        self._drawing_mask = np.zeros((self.SLIDE_HEIGHT, self.SLIDE_WIDTH), dtype=np.uint8)
        
        # Optimize frame processing
        self.PROCESS_WIDTH = 640
        self.PROCESS_HEIGHT = 480
        self.process_matrix = np.array([
            [self.SLIDE_WIDTH/self.PROCESS_WIDTH, 0],
            [0, self.SLIDE_HEIGHT/self.PROCESS_HEIGHT]
        ])
        
        
    def _detect_gesture(self, hand_landmarks):
        """Detect gestures based on hand landmarks."""
        def finger_is_extended(tip_id, pip_id):
            """Check if a finger is extended."""
            return hand_landmarks.landmark[tip_id].y < hand_landmarks.landmark[pip_id].y

        thumb_extended = hand_landmarks.landmark[self.mp_hands.HandLandmark.THUMB_TIP].x < \
                        hand_landmarks.landmark[self.mp_hands.HandLandmark.THUMB_IP].x

        index_extended = finger_is_extended(self.mp_hands.HandLandmark.INDEX_FINGER_TIP,
                                            self.mp_hands.HandLandmark.INDEX_FINGER_PIP)
        middle_extended = finger_is_extended(self.mp_hands.HandLandmark.MIDDLE_FINGER_TIP,
                                            self.mp_hands.HandLandmark.MIDDLE_FINGER_PIP)
        ring_extended = finger_is_extended(self.mp_hands.HandLandmark.RING_FINGER_TIP,
                                            self.mp_hands.HandLandmark.RING_FINGER_PIP)
        pinky_extended = finger_is_extended(self.mp_hands.HandLandmark.PINKY_TIP,
                                            self.mp_hands.HandLandmark.PINKY_PIP)

        if thumb_extended and index_extended and middle_extended and not ring_extended and not pinky_extended:
            return "NEXT"
        elif thumb_extended and index_extended and not middle_extended and not ring_extended and not pinky_extended:
            return "PREVIOUS"
        elif index_extended and middle_extended and not thumb_extended and not ring_extended and not pinky_extended:
            return "DRAW"
        elif index_extended and not thumb_extended and not middle_extended and not ring_extended and not pinky_extended:
            return "POINTER"
        elif thumb_extended and index_extended and middle_extended and ring_extended and pinky_extended:
            return "CLEAR"

        return "NONE"
    
    
    def _handle_gesture(self, gesture, x, y, current_drawing, display):
        """Handle different gestures and their corresponding actions."""
        if gesture != self.last_gesture:
            self.gesture_cooldown = 5
            self.gesture_active = False
            if gesture in ["NEXT", "PREVIOUS", "CLEAR"]:
                self.prev_x = self.prev_y = None

        if self.gesture_cooldown > 0:
            self.gesture_cooldown -= 1

        if self.navigation_timer > 0:
            self.navigation_timer -= 1

        if self.gesture_cooldown == 0:
            if gesture == "DRAW":
                if y > self.CONTROL_HEIGHT:  # Only draw below control panel
                    if self.prev_x is None:
                        self.prev_x, self.prev_y = x, y
                    else:
                        self.draw(x, y, current_drawing)

            elif gesture == "POINTER":
                if y > self.CONTROL_HEIGHT:
                    cv2.circle(display, (x, y), 5, self.pointer_color, -1)
                self.prev_x = self.prev_y = None

            elif not self.is_whiteboard and gesture in ["NEXT", "PREVIOUS"] and self.navigation_timer == 0:
                if not self.gesture_active:
                    if gesture == "NEXT" and self.current_slide < len(self.slides) - 1:
                        self.current_slide += 1
                        self.navigation_timer = 15
                    elif gesture == "PREVIOUS" and self.current_slide > 0:
                        self.current_slide -= 1
                        self.navigation_timer = 15
                    self.gesture_active = True

            elif gesture == "CLEAR":
                if not self.gesture_active:
                    self.clear_drawings()
                    self.gesture_active = True

        self.last_gesture = gesture
        
        
        
    

    def setup_ui(self):
        # Calculate total width of all buttons
        total_buttons = 5  # number of tool buttons
        total_colors = len(self.colors)  # number of colors
        total_sizes = 3  # number of brush sizes
        
        total_width = (total_buttons * self.BUTTON_WIDTH) + \
                    (total_colors * 30) + \
                    (total_sizes * 30) + \
                    (total_buttons + total_colors + total_sizes - 1) * self.BUTTON_SPACING
        
        # Calculate starting x position to center the control panel
        start_x = (self.SLIDE_WIDTH - total_width) // 2
        y = 5
        x = start_x
        
        # Tool buttons
        tools = [
            ("brush", (100, 100, 100), "Brush"),
            ("eraser", (100, 100, 100), "Eraser"),
            ("style", (100, 100, 100), "Style"),
            ("clear", (100, 100, 100), "Clear"),
            ("whiteboard", (100, 100, 100), "Board")
        ]
        
        for name, color, label in tools:
            self.ui_elements[name] = {
                'x': x,
                'y': y,
                'width': self.BUTTON_WIDTH,
                'height': self.CONTROL_HEIGHT - 10,
                'color': color,
                'label': label,
                'active': False
            }
            x += self.BUTTON_WIDTH + self.BUTTON_SPACING
        
        # Color palette
        for color_name, color_value in self.colors.items():
            self.ui_elements[f"color_{color_name}"] = {
                'x': x,
                'y': y,
                'width': 30,
                'height': self.CONTROL_HEIGHT - 10,
                'color': color_value,
                'label': "",
                'active': False
            }
            x += 35
        
        # Brush sizes
        sizes = [(2, "S"), (4, "M"), (6, "L")]
        for size, label in sizes:
            self.ui_elements[f"size_{size}"] = {
                'x': x,
                'y': y,
                'width': 30,
                'height': self.CONTROL_HEIGHT - 10,
                'color': (100, 100, 100),
                'label': label,
                'size': size,
                'active': False
            }
            x += 35

    def draw_ui(self, display):
        # Draw control panel background - now centered
        panel_width = self.SLIDE_WIDTH - 200  # Reduce width to center
        panel_start_x = 100  # Start 100 pixels from left
        cv2.rectangle(display, 
                    (panel_start_x, 0), 
                    (panel_start_x + panel_width, self.CONTROL_HEIGHT), 
                    (50, 50, 50), -1)
        
        # Draw all UI elements
        for name, element in self.ui_elements.items():
            color = element['color']
            if (name == "brush" and self.current_tool == "brush") or \
            (name == "eraser" and self.is_eraser) or \
            (name == "whiteboard" and self.is_whiteboard):
                color = (100, 100, 255)
            
            cv2.rectangle(display, 
                        (element['x'], element['y']), 
                        (element['x'] + element['width'], 
                        element['y'] + element['height']), 
                        color, -1)
            
            if element['label']:
                cv2.putText(display, element['label'], 
                        (element['x'] + 5, element['y'] + 20),
                        cv2.FONT_HERSHEY_SIMPLEX, 0.5, (255, 255, 255), 1)

    def handle_ui_interaction(self, x, y):
        if y > self.CONTROL_HEIGHT:
            # Reset hover states when not in UI area
            for state in self.hover_states.values():
                state['is_hovering'] = False
            return False
            
        for name, element in self.ui_elements.items():
            if (element['x'] <= x <= element['x'] + element['width'] and
                element['y'] <= y <= element['y'] + element['height']):
                
                if name == "brush":
                    self.current_tool = "brush"
                    self.is_eraser = False
                elif name == "eraser":
                    # Handle eraser hover state
                    if not self.hover_states['eraser']['is_hovering'] and self.hover_states['eraser']['cooldown'] == 0:
                        self.is_eraser = not self.is_eraser
                        self.hover_states['eraser']['is_hovering'] = True
                        self.hover_states['eraser']['cooldown'] = self.HOVER_COOLDOWN
                elif name == "whiteboard":
                    # Handle whiteboard hover state
                    if not self.hover_states['whiteboard']['is_hovering'] and self.hover_states['whiteboard']['cooldown'] == 0:
                        self.is_whiteboard = not self.is_whiteboard
                        self.hover_states['whiteboard']['is_hovering'] = True
                        self.hover_states['whiteboard']['cooldown'] = self.HOVER_COOLDOWN
                elif name == "clear":
                    if self.is_whiteboard:
                        self.whiteboard_drawing.fill(0)
                    else:
                        self.drawings[self.current_slide].fill(0)
                elif name.startswith("color_"):
                    color_name = name.split("_")[1]
                    self.current_color = self.colors[color_name]
                elif name.startswith("size_"):
                    self.brush_thickness = element['size']
                elif name == "style":
                    self._cycle_brush_style()
                    
                return True
                
        # Reset hover states when not hovering over eraser or whiteboard
        for state in self.hover_states.values():
            state['is_hovering'] = False
        return False

    def _cycle_brush_style(self):
        current_index = self._brush_styles.index(self.current_brush_style)
        next_index = (current_index + 1) % len(self._brush_styles)
        self.current_brush_style = self._brush_styles[next_index]

    def draw(self, x, y, drawing_surface):
        if self.prev_x is None:
            self.prev_x, self.prev_y = x, y
            return

        if self.is_eraser:
            cv2.circle(drawing_surface, (x, y), self.brush_thickness * 3, (0, 0, 0), -1)
        else:
            if self.current_brush_style == "spray":
                self._spray_paint(x, y, drawing_surface)
            elif self.current_brush_style == "calligraphy":
                self._calligraphy_stroke((self.prev_x, self.prev_y), (x, y), drawing_surface)
            elif self.current_brush_style == "neon":
                self._neon_stroke((self.prev_x, self.prev_y), (x, y), drawing_surface)
            else:  # normal brush
                cv2.line(drawing_surface, (self.prev_x, self.prev_y), (x, y),
                        self.current_color, self.brush_thickness)

        self.prev_x, self.prev_y = x, y

    def _spray_paint(self, x, y, surface):
        for _ in range(20):
            angle = np.random.random() * 360
            radius = np.random.random() * self.brush_thickness * 2
            px = int(x + radius * np.cos(angle))
            py = int(y + radius * np.sin(angle))
            cv2.circle(surface, (px, py), 1, self.current_color, -1)

    def _calligraphy_stroke(self, start_pos, end_pos, surface):
        angle = np.arctan2(end_pos[1]-start_pos[1], end_pos[0]-start_pos[0]) + np.pi/4
        points = []
        width = self.brush_thickness * 2
        
        for i in [-1, 1]:
            dx = width * np.cos(angle) * i
            dy = width * np.sin(angle) * i
            points.append((int(start_pos[0] + dx), int(start_pos[1] + dy)))
            points.append((int(end_pos[0] + dx), int(end_pos[1] + dy)))
        
        pts = np.array(points, np.int32)
        cv2.fillPoly(surface, [pts], self.current_color)

    def _neon_stroke(self, start_pos, end_pos, surface):
        for thickness in range(self.brush_thickness * 2, self.brush_thickness - 1, -1):
            alpha = (thickness - self.brush_thickness) / (self.brush_thickness * 2)
            color = tuple(int(c * alpha) for c in self.current_color)
            cv2.line(surface, start_pos, end_pos, color, thickness)
        cv2.line(surface, start_pos, end_pos, self.current_color, self.brush_thickness)

    def process_frame(self, frame_data: str) -> dict[str, Any] | None:
        try:
            # Cache frame data for repeated processing
            cache_key = hash(frame_data)
            if cache_key in self._frame_cache:
                return self._frame_cache[cache_key]
            
            frame = self._decode_base64_frame(frame_data)
            if frame is None:
                return None

            # Reduce frame resolution for processing
            frame = cv2.resize(frame, (self.PROCESS_WIDTH, self.PROCESS_HEIGHT))  # Process at lower resolution
            
            # Calculate FPS 
            current_time = time.time()
            fps = 1 / (current_time - self.last_frame_time)
            self.fps_history.append(fps)
            self.last_frame_time = current_time

            frame = cv2.flip(frame, 1)
            
            # Convert to RGB more efficiently
            frame_rgb = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
            results = self.hands.process(frame_rgb)

            # Get current display surface
            if self.is_whiteboard:
                display = self.whiteboard
                current_drawing = self.whiteboard_drawing
            else:
                display = self.slides[self.current_slide]
                current_drawing = self.drawings[self.current_slide]

            # Create display copy only if needed
            display_copy = None

            if results.multi_hand_landmarks:
                hand_landmarks = results.multi_hand_landmarks[0]
                
                # Scale coordinates from processed frame to display resolution
                x = int(hand_landmarks.landmark[self.mp_hands.HandLandmark.INDEX_FINGER_TIP].x * self.SLIDE_WIDTH)
                y = int(hand_landmarks.landmark[self.mp_hands.HandLandmark.INDEX_FINGER_TIP].y * self.SLIDE_HEIGHT)

                if not self.handle_ui_interaction(x, y):
                    gesture = self._detect_gesture(hand_landmarks)
                    
                    # Create display copy only when needed
                    if display_copy is None:
                        display_copy = display.copy()
                        
                    self._handle_gesture(gesture, x, y, current_drawing, display_copy)

            # Optimize drawing overlay
            if np.any(current_drawing):  # Only process if there are drawings
                if display_copy is None:
                    display_copy = display.copy()
                mask = cv2.cvtColor(current_drawing, cv2.COLOR_BGR2GRAY)
                display_copy[mask != 0] = current_drawing[mask != 0]

            final_display = display_copy if display_copy is not None else display
            
            # Draw UI
            self.draw_ui(final_display)

            # Optimize status text rendering
            if len(self.fps_history) > 0:  # Only calculate if we have FPS data
                avg_fps = int(sum(self.fps_history) / len(self.fps_history))
                mode_text = "Whiteboard" if self.is_whiteboard else f"Slide {self.current_slide + 1}/{len(self.slides)}"
                status_text = f"{mode_text} - FPS: {avg_fps}"
                cv2.putText(final_display, status_text, (10, self.SLIDE_HEIGHT - 20),
                        cv2.FONT_HERSHEY_SIMPLEX, 0.7, (255, 255, 255), 2)

            # Update cooldown states more efficiently
            for state in self.hover_states.values():
                if state['cooldown'] > 0:
                    state['cooldown'] -= 1

            # Use JPEG encoding with optimized quality for better performance
            result = {
                'frame': f'data:image/jpeg;base64,{self._encode_frame_to_base64(final_display)}',
                'currentSlide': self.current_slide,
                'totalSlides': len(self.slides)
            }
            
            # Cache the result
            if len(self._frame_cache) > self.CACHE_SIZE:
                self._frame_cache.clear()
            self._frame_cache[cache_key] = result
            return result

        except Exception as e:
            logging.error(f"Error processing frame: {e}")
            return None

    def _decode_base64_frame(self, frame_data):
        try:
            # More efficient decoding
            img_bytes = base64.b64decode(frame_data.split(',')[1])
            img_arr = np.frombuffer(img_bytes, np.uint8)
            return cv2.imdecode(img_arr, cv2.IMREAD_REDUCED_COLOR_2)  # Load at reduced resolution
        except Exception as e:
            logging.error(f"Error decoding frame: {e}")
            return None

    def _encode_frame_to_base64(self, frame):
        try:
            # Optimize JPEG encoding for better performance
            encode_param = [int(cv2.IMWRITE_JPEG_QUALITY), 85]  # Balanced quality
            _, buffer = cv2.imencode('.jpg', frame, encode_param)
            return base64.b64encode(buffer).decode('utf-8')
        except Exception as e:
            logging.error(f"Error encoding frame: {e}")
            return None

    def load_ppt(self, ppt_data):
        try:
            # Create temporary files more efficiently
            with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as pptx_temp:
                pptx_temp.write(ppt_data)
                pptx_path = pptx_temp.name

            pdf_path = pptx_path.replace('.pptx', '.pdf')

            # Convert PPTX to PDF
            subprocess.run([
                'C:\\Program Files\\LibreOffice\\program\\soffice.exe',
                '--headless',
                '--convert-to', 'pdf',
                '--outdir', os.path.dirname(pdf_path),
                pptx_path
            ], capture_output=True)  # Capture output for better performance

            # Load PDF and convert slides
            pdf_document = fitz.open(pdf_path)
            
            # Pre-allocate lists for better performance
            self.slides = deque(maxlen=100)
            self.drawings = deque(maxlen=100)
            
            for page_num in range(len(pdf_document)):
                page = pdf_document[page_num]
                pix = page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5))  # Reduced resolution
                
                # More efficient image conversion
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                np_img = np.array(img)
                
                # Optimize resize operation
                slide = cv2.resize(np_img, (self.SLIDE_WIDTH, self.SLIDE_HEIGHT), 
                                 interpolation=cv2.INTER_LINEAR)  # Faster interpolation
                
                slide = cv2.cvtColor(slide, cv2.COLOR_RGB2BGR)
                
                self.slides.append(slide)
                self.drawings.append(np.zeros((self.SLIDE_HEIGHT, self.SLIDE_WIDTH, 3), 
                                           dtype=np.uint8))

            # Cleanup
            pdf_document.close()
            os.unlink(pptx_path)  # More efficient file removal
            os.unlink(pdf_path)
            
            self.current_slide = 0
            return True

        except Exception as e:
            logging.error(f"Error loading PowerPoint: {e}")
            return False

    def next_slide(self):
        if self.current_slide < len(self.slides) - 1:
            self.current_slide += 1
            return True
        return False

    def previous_slide(self):
        if self.current_slide > 0:
            self.current_slide -= 1
            return True
        return False

    def clear_drawings(self):
        if self.is_whiteboard:
            self.whiteboard_drawing.fill(0)
        else:
            self.drawings[self.current_slide].fill(0)

    def toggle_whiteboard(self):
        self.is_whiteboard = not self.is_whiteboard
        
    def cleanup(self):
        self.hands.close()